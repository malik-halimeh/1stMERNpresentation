#!/usr/bin/env python3
"""
build_pptx.py — regenerate the OptiCart presentation decks.

Outputs (next to this script):
  OptiCart_Main.pptx    ~22 slides, STANDARD variant, mirrors presentation_script.html numbering.
  OptiCart_Backup.pptx  ~20 deep-dive slides (B1..B20), mapped from rehearsal.html.

Design system mirrors the HTML kit: dark slate #0F172A, orange #EA580C titles,
teal #0D9488 secondary, violet #8B5CF6 for Mahmod. Key diagrams are drawn as NATIVE
shapes so the decks are fully offline (no image conversion needed).

Run:  python build_pptx.py
It generates both decks, then runs an inspection pass for overflow / empty titles
and prints a report.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import os, math

# ── palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0F, 0x17, 0x2A)
SURFACE = RGBColor(0x1E, 0x29, 0x3B)
SURF2   = RGBColor(0x27, 0x34, 0x49)
TEXT    = RGBColor(0xE2, 0xE8, 0xF0)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
WHITE   = RGBColor(0xF8, 0xFA, 0xFC)
ORANGE  = RGBColor(0xEA, 0x58, 0x0C)
ORANGE2 = RGBColor(0xFB, 0x92, 0x3C)
TEAL    = RGBColor(0x0D, 0x94, 0x88)
TEAL2   = RGBColor(0x2D, 0xD4, 0xBF)
VIOLET  = RGBColor(0x8B, 0x5C, 0xF6)
VIOLET2 = RGBColor(0xA7, 0x8B, 0xFA)
BORDER  = RGBColor(0x33, 0x41, 0x55)
GREEN   = RGBColor(0x34, 0xD3, 0x99)
RED     = RGBColor(0xF8, 0x71, 0x71)

PRESENTER = {
    "Haya":   (TEAL,   TEAL2),
    "Malik":  (ORANGE, ORANGE2),
    "Mahmod": (VIOLET, VIOLET2),
    "All":    (SURF2,  MUTED),
}

EMU = 914400
SW, SH = 13.333, 7.5   # 16:9 inches

# ── low-level helpers ─────────────────────────────────────────────────────────

def _solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color

def _noline(shape):
    shape.line.fill.background()

def _line(shape, color, w=1.0):
    shape.line.color.rgb = color; shape.line.width = Pt(w)

def bg(slide, color=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    _solid(r, color); _noline(r)
    slide.shapes._spTree.remove(r._element); slide.shapes._spTree.insert(2, r._element)
    return r

def textbox(slide, l, t, w, h, runs, size=16, color=TEXT, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI",
            line_spacing=1.12, space_after=6):
    """runs: a string, or a list of (text,size,color,bold) paragraphs, or list of
    (text,size,color,bold,level) for bullets."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    if isinstance(runs, str):
        runs = [(runs, size, color, bold)]
    for i, item in enumerate(runs):
        text, sz, col, bd = item[0], item[1], item[2], item[3]
        lvl = item[4] if len(item) > 4 else 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.level = lvl
        p.line_spacing = line_spacing; p.space_after = Pt(space_after)
        r = p.add_run(); r.text = text
        r.font.size = Pt(sz); r.font.bold = bd; r.font.color.rgb = col; r.font.name = font
    return tb

def box(slide, l, t, w, h, fill, line=None, text=None, tcolor=WHITE, tsize=13,
        bold=True, radius=True, line_w=1.25, align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    _solid(shp, fill)
    if line: _line(shp, line, line_w)
    else: _noline(shp)
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    if text is not None:
        for j, ln in enumerate(text.split("\n")):
            pp = p if j == 0 else tf.add_paragraph()
            pp.alignment = align
            rr = pp.add_run(); rr.text = ln
            rr.font.size = Pt(tsize); rr.font.bold = bold
            rr.font.color.rgb = tcolor; rr.font.name = "Segoe UI"
    return shp

def connector(slide, x1, y1, x2, y2, color=MUTED, w=1.25, arrow=True, dashed=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w)
    try:
        ln = cn.line._get_or_add_ln()
        if dashed:
            d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'}); ln.append(d)
        if arrow:
            te = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
            ln.append(te)
    except Exception:
        pass
    return cn

# ── slide scaffolding ─────────────────────────────────────────────────────────

def base(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg(s)
    return s

def header(slide, num, title, presenter="All", accent=ORANGE):
    # left accent band coloured by presenter section
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), Inches(SH))
    _solid(band, accent); _noline(band)
    tsize = 24 if len(str(title)) > 32 else 30
    textbox(slide, 0.55, 0.36, 9.9, 1.12, [(title, tsize, WHITE, True)],
            anchor=MSO_ANCHOR.MIDDLE)
    # accent underline
    u = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.4),
                               Inches(2.4), Inches(0.05))
    _solid(u, accent); _noline(u)
    # presenter chip
    pc, pt = PRESENTER.get(presenter, PRESENTER["All"])
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(SW-2.55), Inches(0.5), Inches(2.0), Inches(0.42))
    _solid(chip, SURFACE); _line(chip, pc, 1.25)
    cf = chip.text_frame; cf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = cf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = presenter; cr.font.size = Pt(13); cr.font.bold = True
    cr.font.color.rgb = pt; cr.font.name = "Segoe UI"
    # slide number + brand
    textbox(slide, SW-2.55, 0.98, 2.0, 0.3, [("OptiCart Appliances", 9, MUTED, False)],
            align=PP_ALIGN.CENTER)
    textbox(slide, SW-1.2, SH-0.5, 0.9, 0.3, [(str(num), 11, MUTED, False)],
            align=PP_ALIGN.RIGHT)

def bullets(slide, items, l=0.7, t=1.7, w=8.4, h=5.0, size=17, gap=8):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        prefix = "•  " if lvl == 0 else "–  "
        runs.append((prefix + txt, size if lvl == 0 else size-2,
                     TEXT if lvl == 0 else MUTED, False, lvl))
    textbox(slide, l, t, w, h, runs, space_after=gap, line_spacing=1.14)

def note_card(slide, l, t, w, h, title, body, accent=TEAL):
    c = box(slide, l, t, w, h, SURFACE, line=accent, text=None, radius=True, line_w=1.5)
    textbox(slide, l+0.2, t+0.14, w-0.4, h-0.3,
            [(title, 14, WHITE, True)] + [(body, 12.5, MUTED, False)],
            space_after=4, line_spacing=1.12)

# ══════════════════════════════════════════════════════════════════════════════
#  MAIN DECK
# ══════════════════════════════════════════════════════════════════════════════

def build_main(path):
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)

    # 1 · Title
    s = base(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.55), Inches(SW), Inches(0.06))
    _solid(bar, ORANGE); _noline(bar)
    textbox(s, 0.8, 1.4, 11.7, 1.2, [("OptiCart Appliances", 46, WHITE, True)])
    textbox(s, 0.8, 2.72, 11.7, 0.9,
            [("A full-stack MERN e-commerce platform — storefront + role-scoped back office", 18, ORANGE2, False)])
    textbox(s, 0.8, 3.7, 11.7, 1.6, [
        ("Haya  ·  Frontend  (React / Vite / TypeScript)", 17, TEAL2, True),
        ("Malik  ·  Backend & Architecture  (Express / TypeScript)", 17, ORANGE2, True),
        ("Mahmod  ·  Database  (MongoDB / Mongoose)", 17, VIOLET2, True),
    ], space_after=8)
    textbox(s, 0.8, SH-0.7, 11.7, 0.4, [("Get the money and the stock right — every time.", 14, MUTED, False)])

    # 2 · The Idea
    s = base(prs); header(s, 2, "The Idea", "Malik")
    bullets(s, [
        "Buying an appliance online is a trust exercise: accurate stock, a safe checkout, promises kept after the sale.",
        "So we built both sides of one system — a storefront to shop and a back office to run the business.",
        "Three roles: customer, inventory manager, super admin.",
        "North star: get the money and the stock exactly right, every single time.",
    ])

    # 3 · Agenda
    s = base(prs); header(s, 3, "Agenda", "Malik")
    bullets(s, [
        "Planning — requirements, data design, architecture.",
        "Live demo — each of us presents the part we built.",
        "What we learned — and what we'd still improve.",
        "Future feature — the Smart Energy Advisor.",
        "Q&A — jump in any time.",
    ])

    # 4 · Requirements & Roles
    s = base(prs); header(s, 4, "Planning — Requirements & Roles", "Mahmod")
    note_card(s, 0.7, 1.7, 5.7, 2.4, "Functional",
              "Browse · search · filter · guest cart & wishlist · checkout · orders · verified reviews. "
              "Admin: products, categories, orders, coupons, low-stock, procurement, users, analytics, audit.", TEAL)
    note_card(s, 6.9, 1.7, 5.7, 2.4, "Non-functional",
              "Security (token auth + RBAC) · correctness under load (no overselling, cents money) · "
              "responsiveness (mobile) · maintainability (typed end to end).", ORANGE)
    box(s, 0.7, 4.4, 11.9, 1.9, SURFACE, line=BORDER, radius=True, line_w=1.25, text=None)
    textbox(s, 1.0, 4.55, 11.4, 1.7, [
        ("Roles are supersets, enforced on the server:", 15, WHITE, True),
        ("customer  ⊂  inventory_manager  ⊂  super_admin", 18, ORANGE2, True),
        ("authenticate → req.user, then authorize(...roles) gates each route (middleware/auth.ts).", 13, MUTED, False),
    ], space_after=6)

    # 5 · Design Direction
    s = base(prs); header(s, 5, "Planning — Design Direction", "Haya")
    bullets(s, [
        "One visual language: a Tailwind token system (colors, spacing, type) defined once, consumed everywhere.",
        "That's why storefront and admin feel like the same product.",
        "Core flows sketched first: product list + filters, detail + variant picker, cart → checkout.",
        "The rule: never a frozen screen — every load has a skeleton, every action has feedback.",
    ])

    # 6 · ERD (native diagram)
    s = base(prs); header(s, 6, "Planning — Database Design (ERD)", "Mahmod")
    _erd(s)
    textbox(s, 0.7, 6.55, 11.9, 0.8, [
        ("Embed what's owned & read together (variants in Product); reference what's shared (Order → Product). "
         "Orders/Purchases store snapshots so history never changes. Money = integer cents. 14 collections.", 12.5, MUTED, False)])

    # 7 · Architecture (native diagram)
    s = base(prs); header(s, 7, "Architecture & Request Lifecycle", "Malik")
    _arch(s)
    textbox(s, 0.7, 6.5, 11.9, 0.85, [
        ("Every request: helmet → CORS → (raw webhook body BEFORE json) → rate limit → authN/authZ → controller → "
         "one error handler. Raw-before-JSON is what makes Stripe payments tamper-proof.", 12.5, MUTED, False)])

    # 8 · Sprint & Pre-Dev
    s = base(prs); header(s, 8, "Sprint Board & Pre-Dev Decisions", "Malik")
    bullets(s, [
        "Planned as 11 epics → 40 sub-tasks → 311 story points, split evenly by strength (Jira board).",
        "Big bets made up front, so we never relitigated them:",
        ("MERN + TypeScript everywhere; MongoDB document model for variant-heavy products.", 1),
        ("Integer cents for money; two-token auth; Stripe; Tailwind token system.", 1),
        "That's why the codebase is coherent, not stitched together.",
    ])

    # 9 · Demo Handoff
    s = base(prs); header(s, 9, "Live Demo", "Malik")
    textbox(s, 0.7, 2.6, 11.9, 2.0, [
        ("One shopper's journey, end to end.", 30, WHITE, True),
        ("Each of us jumps in on the part we built — Haya shops, Malik secures & sells, Mahmod proves the data.", 17, ORANGE2, False),
    ], space_after=14)

    # 10 · Haya demo
    s = base(prs); header(s, 10, "Demo — Frontend", "Haya", accent=TEAL)
    bullets(s, [
        "Guest browsing — live category nav (by slug), URL-driven search (shareable, back-safe).",
        "Filters update list + URL together; skeletons on every load.",
        "Product detail — variant picker updates price, gallery, stock; add-to-cart targets the SKU.",
        "Wishlist + cart (with toast) as a guest — cart lives in localStorage.",
        "Start checkout → the moment I sign in, the guest cart merges. Handoff to Malik.",
    ])

    # 11 · Malik auth
    s = base(prs); header(s, 11, "Demo — Auth & Security", "Malik", accent=ORANGE)
    bullets(s, [
        "Sign in → guest cart merges onto the account.",
        "Network tab: access token in the Authorization header — memory only, never localStorage.",
        "On 401 the client silently refreshes (http-only cookie) and retries — rotated, with theft detection.",
        "Hit an admin endpoint as a customer → 403. RBAC enforced on the server, not hidden in the UI.",
    ])

    # 12 · Malik checkout
    s = base(prs); header(s, 12, "Demo — Checkout & Order Lifecycle", "Malik", accent=ORANGE)
    bullets(s, [
        "Apply a coupon, pay: the server re-prices the cart and re-validates the coupon from scratch.",
        "Stock decrements only on payment success — an atomic guarded op; two buyers can't both win.",
        "Admin queue: confirm → ship → deliver, each transition validated (no skipping).",
        "Each step fires the customer's bell + email. Payment is machine; fulfilment stays human.",
    ])

    # 13 · Mahmod integrity
    s = base(prs); header(s, 13, "Demo — Data Integrity in Action", "Mahmod", accent=VIOLET)
    bullets(s, [
        "Delivered order unlocks a verified review — unique index (one per customer/product) + transaction recompute.",
        "Order snapshots: rename the product, the past order is unchanged.",
        "Audit log is append-only — editing an entry is rejected by the model hook.",
        "Information you can quietly change isn't a record — so we made it impossible to change.",
    ])

    # 14 · Procurement
    s = base(prs); header(s, 14, "Procurement & Weighted-Average Cost", "Mahmod", accent=VIOLET)
    bullets(s, [
        "Record a stock purchase: stock rises AND weighted-average cost blends old + new by quantity.",
        "Both happen in a single atomic database operation — stock and cost can never disagree.",
        "The dashboard month-spend tile moves in real time.",
        "If the record fails after the stock update, we auto-undo it — a compensating step.",
    ])

    # 15 · Analytics & Audit
    s = base(prs); header(s, 15, "Analytics & Audit Trail", "Malik")
    bullets(s, [
        "Super admin sees revenue, order volume, top products, customer growth — and now purchase spend.",
        "Revenue vs cost in one place — real MongoDB aggregation pipelines, not client-side math.",
        "Under everything: an immutable audit log — who, what, before & after — deliberately indexed.",
    ])

    # 16 · Learned
    s = base(prs); header(s, 16, "What We Learned", "All")
    note_card(s, 0.7, 1.8, 3.9, 3.9, "Haya — Single-flight refresh",
              "N concurrent 401s trigger ONE refresh, not N. The fix wasn't more state — it was one point "
              "of coordination. And why a token belongs in memory, not storage.", TEAL)
    note_card(s, 4.75, 1.8, 3.9, 3.9, "Malik — Compensating transactions",
              "When an operation spans two documents, make the second step reversible and undo the first on "
              "failure — instead of heavy machinery. It made our costing feature safe.", ORANGE)
    note_card(s, 8.8, 1.8, 3.8, 3.9, "Mahmod — Partial indexes",
              "A unique constraint on a filtered subset turned a real bug into an elegant rule: one active "
              "low-stock alert per variant, unlimited history — enforced by the database.", VIOLET)

    # 17 · Challenging
    s = base(prs); header(s, 17, "Still Challenging (honest)", "All")
    bullets(s, [
        "Haya — automated frontend testing: we have a backend suite; UI component tests are next.",
        "Malik — designing for horizontal scale: Redis + a cron leader-lock, not yet run multi-instance.",
        "Mahmod — a formal data-migration workflow for evolving live data safely.",
        "None of these block the product — they're the edges we're growing into.",
    ])

    # 18 · Evaluation mapping
    s = base(prs); header(s, 18, "Evaluation Criteria — How We Meet Them", "Malik")
    _eval_table(s)

    # 19 · Security found/fixed/proven
    s = base(prs); header(s, 19, "Security — Found → Fixed → Proven", "Malik")
    _fixed_table(s)

    # 20 · Honesty
    s = base(prs); header(s, 20, "Honesty — Open Items & Plan", "Malik")
    bullets(s, [
        "Rate limiting + scheduler assume a single instance → Redis store + cron leader-lock.",
        "Dev-format logging; no CI/Docker yet; client bundle could be code-split.",
        "None is a happy-path correctness bug — they're scaling & hygiene items, each with a clear fix.",
        "We'd rather raise them ourselves than have you find them.",
    ])

    # 21 · Future — Smart Energy Advisor
    s = base(prs); header(s, 21, "Future — Smart Energy Advisor", "All")
    bullets(s, [
        "Two $899 fridges look identical — but one can cost double to run.",
        "Store each appliance's energy use; show its real yearly running cost under Lebanese conditions.",
        "EDL grid tariff + private generator rate — and sort the whole catalog by total cost of ownership.",
    ], t=1.6, h=2.6)
    box(s, 0.7, 4.15, 11.9, 0.95, SURFACE, line=ORANGE, radius=True, line_w=1.5,
        text="yearlyCost = kWh/yr × (gridShare×EDL + genShare×generator)     5-yr TCO = price + 5 × yearlyCost",
        tcolor=ORANGE2, tsize=14, bold=True)
    box(s, 0.7, 5.35, 11.9, 1.05, RGBColor(0x2A,0x14,0x08), line=ORANGE, radius=True, line_w=1.25,
        text="\"This fridge costs $120/year to run — this one $55.\"  In Lebanon, that's the sentence that closes the sale.",
        tcolor=WHITE, tsize=15, bold=True)

    # 22 · Thank you
    s = base(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), Inches(SW), Inches(0.06))
    _solid(bar, ORANGE); _noline(bar)
    textbox(s, 0.8, 1.5, 11.7, 1.0, [("Thank you", 46, WHITE, True)])
    textbox(s, 0.8, 2.7, 11.7, 0.9, [("Questions? We know exactly where everything lives — name it and we'll open the file.", 17, ORANGE2, False)])
    textbox(s, 0.8, 3.7, 11.7, 1.4, [
        ("Haya · Frontend", 18, TEAL2, True),
        ("Malik · Backend & Architecture", 18, ORANGE2, True),
        ("Mahmod · Database", 18, VIOLET2, True),
    ], space_after=8)

    prs.save(path)
    return prs

# ── native diagrams ───────────────────────────────────────────────────────────

def _arch(s):
    # three bands
    y = 1.65
    box(s, 0.7, y, 11.9, 1.15, RGBColor(0x2A,0x14,0x08), line=ORANGE, radius=True, line_w=1.25, text=None)
    textbox(s, 0.85, y+0.06, 4, 0.3, [("CLIENT · React 19 SPA", 11, ORANGE2, True)])
    box(s, 1.0, y+0.42, 3.3, 0.6, SURFACE, line=ORANGE, text="Pages + UI + Contexts", tsize=12)
    box(s, 4.7, y+0.42, 3.3, 0.6, SURFACE, line=ORANGE, text="services/api.ts (Axios+refresh)", tsize=11)
    box(s, 8.4, y+0.42, 3.9, 0.6, SURFACE, line=ORANGE, text="single-flight silent refresh", tsize=11)

    y2 = 3.15
    box(s, 0.7, y2, 11.9, 1.35, RGBColor(0x07,0x22,0x20), line=TEAL, radius=True, line_w=1.25, text=None)
    textbox(s, 0.85, y2+0.06, 4, 0.3, [("SERVER · Express + TypeScript", 11, TEAL2, True)])
    for i,(lbl,ww) in enumerate([("middleware\nchain",2.5),("16 routers",2.0),("controllers",2.3),("one errorHandler",3.0)]):
        xs = 1.0 + sum([2.5,2.0,2.3][:i]) + i*0.35
        box(s, xs, y2+0.5, ww, 0.62, SURFACE, line=TEAL, text=lbl, tsize=11)

    y3 = 4.95
    box(s, 0.7, y3, 11.9, 1.15, RGBColor(0x1A,0x12,0x2E), line=VIOLET, radius=True, line_w=1.25, text=None)
    textbox(s, 0.85, y3+0.06, 4, 0.3, [("DATA & EXTERNAL", 11, VIOLET2, True)])
    for i,lbl in enumerate(["MongoDB · 14 col.","Stripe","Cloudinary","SMTP / console"]):
        box(s, 1.0+i*2.95, y3+0.42, 2.6, 0.6, SURFACE, line=VIOLET, text=lbl, tsize=11)
    # connectors between bands
    connector(s, 6.65, y+1.02, 6.65, y2, TEAL2)
    connector(s, 6.65, y2+1.35, 6.65, y3, VIOLET2)

def _erd(s):
    # hubs
    box(s, 5.3, 1.65, 2.7, 0.72, RGBColor(0x2A,0x14,0x08), line=ORANGE, text="USER", tsize=15)
    box(s, 5.3, 4.55, 2.7, 0.72, RGBColor(0x1A,0x12,0x2E), line=VIOLET, text="PRODUCT", tsize=15)
    # user satellites (left)
    left = [("ORDER",1.9),("CART / WISHLIST",1.1),("REVIEW",3.1)]
    for lbl,ty in left:
        box(s, 0.7, ty, 2.7, 0.55, SURFACE, line=TEAL, text=lbl, tsize=11.5)
        connector(s, 3.4, ty+0.27, 5.3, 2.0, MUTED, arrow=True)
    # product satellites (right)
    right = [("REVIEW",1.7),("LOWSTOCKALERT",2.5),("PURCHASE",3.3),("PRODUCTEVENT (TTL)",4.1),("CATEGORY",4.9)]
    for lbl,ty in right:
        box(s, 9.9, ty, 2.7, 0.5, SURFACE, line=VIOLET, text=lbl, tsize=10.5)
        connector(s, 9.9, ty+0.25, 8.0, 4.9, MUTED, arrow=True)
    # user -> product (places / order)
    connector(s, 6.65, 2.37, 6.65, 4.55, ORANGE, arrow=True)
    textbox(s, 6.75, 3.25, 2.0, 0.3, [("orders", 10, ORANGE2, False)])

def _eval_table(s):
    rows = [
        ("Criterion", "Evidence (file)", "Who"),
        ("Backend architecture", "middleware chain; two-phase pay→confirm (order.ts)", "Malik"),
        ("Security", "rotation+reuse (auth.ts); fail-closed webhook; test_fixes 20/20", "Malik"),
        ("Database design", "snapshots; partial + unique indexes; append-only audit", "Mahmod"),
        ("UI/UX & responsiveness", "token system; live nav; mobile drawer (StorefrontLayout)", "Haya"),
        ("Functionality", "26 verified features — merge, checkout, procurement, analytics", "All"),
        ("Code quality", "one error pipeline; typed schemas; regression suite", "Malik/Haya"),
    ]
    _table(s, rows, top=1.75, widths=[3.6, 6.5, 1.8], fsize=12)

def _fixed_table(s):
    rows = [
        ("Found", "Fixed", "Proven"),
        ("Webhook trusted unverified input", "raw body + fail-closed (400)", "test_fixes.ts"),
        ("Oversell race on stock", "atomic guarded $inc", "20 / 20 checks"),
        ("Refunds were status-only", "real Stripe refund + audit", "clean tsc builds"),
        ("JWT secret used a silent fallback", "lazy read + prod fail-fast", "HTTP smoke test"),
    ]
    _table(s, rows, top=1.9, widths=[4.6, 4.6, 2.7], fsize=13)

def _table(s, rows, top, widths, fsize=12):
    left = 0.7; total_h = 0.62 * len(rows)
    gt = s.shapes.add_table(len(rows), len(widths), Inches(left), Inches(top),
                            Inches(sum(widths)), Inches(total_h)).table
    for j, wv in enumerate(widths):
        gt.columns[j].width = Inches(wv)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = SURF2 if i == 0 else (SURFACE if i % 2 else RGBColor(0x22,0x2E,0x44))
            c.margin_left = Inches(0.08); c.margin_right = Inches(0.08)
            c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = val
            r.font.size = Pt(fsize); r.font.name = "Segoe UI"
            r.font.bold = (i == 0)
            r.font.color.rgb = WHITE if i == 0 else TEXT
    return gt

# ══════════════════════════════════════════════════════════════════════════════
#  BACKUP DECK  (B1..B20)
# ══════════════════════════════════════════════════════════════════════════════

def build_backup(path):
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)

    # B1 title
    s = base(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(SW), Inches(0.06))
    _solid(bar, VIOLET); _noline(bar)
    textbox(s, 0.8, 1.5, 11.7, 1.0, [("Backup — Deep Dives", 42, WHITE, True)])
    textbox(s, 0.8, 2.75, 11.7, 0.6, [("Jump here on deep-dive questions (mapping in rehearsal.html).", 17, VIOLET2, False)])

    def bslide(num, title, who="Malik", accent=ORANGE):
        s = base(prs); header(s, "B"+str(num), title, who, accent=accent); return s

    # B2 auth sequence (native)
    s = bslide(2, "Auth — Login, Rotation, Reuse Detection", "Malik")
    _auth_seq(s)

    # B3 middleware chain
    s = bslide(3, "Middleware Chain — Order Matters", "Malik")
    _chain(s)

    # B4 request lifecycle
    s = bslide(4, "Request Lifecycle", "Malik")
    bullets(s, [
        "helmet → CORS(credentials) → raw webhook body (before json!) → express.json + cookies →",
        "morgan → rate limit (1000/15m) → authenticate → authorize(...roles) → controller.",
        "Success → JSON. Any throw → one errorHandler → { error: { message, code } }.",
        "Boot order: env → Mongo(retry) → syncIndexes → middleware → routers → cron → listen.",
    ])

    # B5 ERD
    s = bslide(5, "Data Model — ERD (14 collections)", "Mahmod", accent=VIOLET); _erd(s)

    # B6 embed vs reference
    s = bslide(6, "Embed vs Reference", "Mahmod", accent=VIOLET)
    note_card(s, 0.7, 1.8, 5.8, 3.4, "Embed (owned, read together)",
              "Variants inside Product; addresses inside User; refund subdoc inside Order. One query loads "
              "everything, updates atomically. Threshold: dozens fine, hundreds → split out.", VIOLET)
    note_card(s, 6.8, 1.8, 5.8, 3.4, "Reference (shared, independent)",
              "Order → Product, Review → User+Product, Purchase → Product. Plus snapshots on Order & Purchase "
              "so history is immutable even when the source changes.", TEAL)

    # B7 snapshots & cents
    s = bslide(7, "Snapshots & Money-as-Cents", "Mahmod", accent=VIOLET)
    bullets(s, [
        "Orders freeze item name + price; Purchases freeze product + buyer name.",
        "Rename or reprice a product later → past records never change. History wins over normalization.",
        "Every monetary field is an integer count of cents — floats can't represent money exactly.",
        "Format to dollars only at the very edge (UI); never do arithmetic on displayed dollars.",
    ])

    # B8 index strategy
    s = bslide(8, "Index Strategy — Rules, not just Speed", "Mahmod", accent=VIOLET)
    _table(s, [
        ("Index", "Serves"),
        ("unique (user, product)", "one review per customer per product"),
        ("partial unique (active only)", "one active low-stock alert per variant"),
        ("TTL 180d on ProductEvent", "self-pruning view stream"),
        ("text (name, description)", "catalog search ranking"),
        ("compound on AuditLog", "filter by actor / action / target"),
    ], top=1.9, widths=[5.0, 6.9], fsize=13)

    # B9 costing pipeline
    s = bslide(9, "Weighted-Average Costing Pipeline", "Mahmod", accent=VIOLET)
    box(s, 0.7, 1.8, 11.9, 1.0, SURFACE, line=ORANGE, radius=True, line_w=1.5,
        text="newCost = (oldStock×oldCost + qty×unitCost) ÷ (oldStock + qty)",
        tcolor=ORANGE2, tsize=17, bold=True)
    bullets(s, [
        "One atomic aggregation-pipeline findOneAndUpdate with $map adds stock AND blends cost together.",
        "Stock and cost can never disagree — no window between them.",
        "If the Purchase record fails afterward, a compensating update reverts the change.",
        "Delete mirrors it and refuses (409) if the revert would drive stock negative.",
    ], t=3.0, h=3.2)

    # B10 checkout & webhook
    s = bslide(10, "Checkout & Fail-Closed Webhook", "Malik"); _checkout(s)

    # B11 atomic stock
    s = bslide(11, "Atomic Stock — No Oversell", "Malik")
    box(s, 0.7, 1.9, 11.9, 1.0, SURFACE, line=ORANGE, radius=True, line_w=1.5,
        text="updateOne({ _id, 'variants.sku': sku, 'variants.stock': { $gte: qty } }, { $inc: { 'variants.$.stock': -qty } })",
        tcolor=ORANGE2, tsize=13, bold=True)
    bullets(s, [
        "The guard stock ≥ qty lives in the query — the database arbitrates concurrent orders.",
        "The loser of a race matches zero documents; we clamp to 0 and log a warning.",
        "Stock only decrements on payment_intent.succeeded — a failed payment leaves inventory untouched.",
    ], t=3.1, h=3.0)

    # B12 refund flow
    s = bslide(12, "Refund Flow", "Malik")
    bullets(s, [
        "stripe.refunds.create (mock id if keyless) runs BEFORE the status persists.",
        "The embedded refund subdoc is populated; a dedicated refund_decision audit entry is written.",
        "A Stripe failure aborts with 502 REFUND_FAILED — status never claims money moved when it didn't.",
    ])

    # B13 error pipeline
    s = bslide(13, "Error Handling — One Envelope", "Malik")
    box(s, 0.7, 1.9, 11.9, 0.95, SURFACE, line=ORANGE, radius=True, line_w=1.5,
        text='{ "error": { "message": "Coupon has expired", "code": "COUPON_EXPIRED" } }',
        tcolor=ORANGE2, tsize=15, bold=True)
    bullets(s, [
        "Controllers throw AppError(message, status, code); one errorHandler serialises everything.",
        "Zod validation failures → 400 with field detail; unknown errors → generic 500, no stack leak.",
        "Predictable contract the client can branch on and localise.",
    ], t=3.1, h=3.0)

    # B14 security layers
    s = bslide(14, "Security Layers", "Malik")
    bullets(s, [
        "Transport: helmet headers, CORS allow-list, SameSite=Strict refresh cookie.",
        "Auth: bcrypt(12), 15-min memory JWT, rotating hashed refresh, reuse detection.",
        "Input: Zod allow-lists strip unknown fields (no role/isActive smuggling).",
        "Payments: fail-closed webhook. Secrets: prod boot fails without JWT_ACCESS_SECRET.",
    ])

    # B15 found/fixed/proven
    s = bslide(15, "Found → Fixed → Proven", "Malik"); _fixed_table(s)

    # B16 scalability
    s = bslide(16, "Scalability — What Breaks at 10×", "Malik")
    bullets(s, [
        "In-memory rate limiter splits per process → move to Redis.",
        "In-process node-cron double-runs on 2 instances → add a leader-lock.",
        "Orders, payments, stock, costing are already DB-atomic — they scale as-is.",
        "Also: structured logging, CI + Docker, client code-splitting (546 kB bundle).",
    ])

    # B17 code quality
    s = bslide(17, "Code Quality Evidence", "Malik", accent=TEAL)
    bullets(s, [
        "End-to-end TypeScript — client/server mismatches are compile errors.",
        "One AppError pipeline; Zod validators separated; a reusable UI kit (DataTable, SearchBox ×7 pages).",
        "Regression suite server/src/test_fixes.ts (20/20) + clean builds + HTTP smoke test.",
        "Honest smells owned: server/dist tracked; some 600-line controllers; duplicated low-stock helper.",
    ])

    # B18 energy advisor detail
    s = bslide(18, "Smart Energy Advisor — Feasibility", "All", accent=ORANGE)
    bullets(s, [
        "Schema: an optional energySpec subdoc (sibling of meta) — zero migration, existing docs simply lack it.",
        "Filter/sort: one filter branch (maxYearlyCost) + one sort branch (tco_asc) on the existing pipeline.",
        "Precompute estYearlyCostCents on save — same precedent as Purchase.totalCostCents.",
        "Admin-configurable tariffs (EDL + generator) — never hardcoded; recompute nightly on change.",
    ])

    # B19 git story
    s = bslide(19, "Git — Honest Before / After", "Malik")
    bullets(s, [
        "26 commits; exemplary messages (b9d48b3, de8c9cb) alongside a couple of low-effort ones.",
        "A brief .env leak-and-cleanup cycle — nothing live is exposed now.",
        "server/dist re-tracked (47 files) — the top hygiene fix on our list.",
        "We show this openly: owning history beats pretending it's spotless.",
    ])

    # B20 thanks
    s = base(prs)
    textbox(s, 0.8, 2.6, 11.7, 1.2, [("Deep-dive backup ends here", 34, WHITE, True),
            ("Ask us anything — we'll open the exact file.", 17, VIOLET2, False)], space_after=12)

    prs.save(path)
    return prs

def _auth_seq(s):
    # 3 lifelines
    cols = [(2.2, "Browser\n(memory token)", ORANGE), (6.6, "Express /api/auth", TEAL), (11.0, "MongoDB users", VIOLET)]
    for x, lbl, c in cols:
        box(s, x-1.35, 1.7, 2.7, 0.7, SURFACE, line=c, text=lbl, tsize=11.5)
        connector(s, x, 2.4, x, 6.3, BORDER, arrow=False, w=1.0)
    msgs = [
        (2.75, 2.2, 6.6, "POST /login (email, password)", TEAL2),
        (2.9, 6.6, 11.0, "store SHA-256(refresh)", VIOLET2),
        (2.9, 6.6, 2.2, "JWT 15m + refresh cookie (httpOnly)", ORANGE2),
        (3.55, 2.2, 6.6, "POST /refresh (single-flight)", TEAL2),
        (3.55, 6.6, 11.0, "findOneAndUpdate: new hash, prev, grace 30s", VIOLET2),
    ]
    for y, x1, x2, lbl, col in msgs:
        connector(s, x1, y, x2, y, MUTED, arrow=True)
        textbox(s, min(x1, x2)+0.1, y-0.34, abs(x2-x1)-0.2, 0.34, [(lbl, 10.5, col, False)],
                align=PP_ALIGN.CENTER)
    box(s, 0.9, 4.4, 11.5, 1.75, RGBColor(0x14,0x1B,0x2E), line=BORDER, radius=True, text=None, line_w=1.0)
    textbox(s, 1.1, 4.5, 11.1, 1.6, [
        ("alt — three outcomes on refresh:", 12, MUTED, True),
        ("✔  winner → new access token + rotated cookie", 12.5, GREEN, False),
        ("↺  concurrent within 30s grace → new access token, no re-rotate", 12.5, RGBColor(0xFB,0xBF,0x24), False),
        ("✘  stale outside grace → null ALL hashes → 403 AUTH_SESSION_COMPROMISED", 12.5, RED, False),
    ], space_after=4)

def _chain(s):
    labels = ["helmet", "CORS", "raw webhook\n(before json)", "json + cookies", "rate limit", "authN / authZ", "controller"]
    cols = [TEAL, TEAL, ORANGE, TEAL, TEAL, ORANGE, VIOLET]
    x = 0.7; y = 2.3; w = 1.6; h = 0.95
    xs = []
    for i, (lbl, c) in enumerate(zip(labels, cols)):
        bx = 0.7 + i*1.78
        xs.append(bx)
        box(s, bx, y, w, h, SURFACE, line=c, text=lbl, tsize=11)
        if i > 0:
            connector(s, xs[i-1]+w, y+h/2, bx, y+h/2, MUTED, arrow=True)
    connector(s, xs[-1]+w/2, y+h, xs[-1]+w/2, y+h+0.5, MUTED, arrow=True)
    box(s, xs[-1]-0.4, y+h+0.5, 2.4, 0.7, RGBColor(0x22,0x2E,0x44), line=GREEN, text="JSON response", tsize=12)
    textbox(s, 0.7, 4.3, 11.9, 1.2, [
        ("The raw-webhook mount comes BEFORE express.json so Stripe's signed bytes survive —", 14, WHITE, True),
        ("parse first and the signature can't be verified. Ordering is correctness, not cosmetics.", 14, MUTED, False),
    ], space_after=4)

def _checkout(s):
    steps = [
        ("POST /checkout — re-price cart, re-validate coupon (server-side)", TEAL),
        ("create Stripe PaymentIntent (or mock if keyless)", TEAL),
        ("Stripe → webhook (raw body): verify signature FAIL-CLOSED (400)", ORANGE),
        ("atomic $inc stock guarded by stock ≥ qty (clamp 0 on race)", ORANGE),
        ("create Order (server prices, line snapshots) + purchase event", VIOLET),
    ]
    y = 1.8
    for i, (lbl, c) in enumerate(steps):
        box(s, 1.6, y, 10.1, 0.72, SURFACE, line=c, text=lbl, tsize=12.5, align=PP_ALIGN.LEFT)
        if i < len(steps)-1:
            connector(s, 6.65, y+0.72, 6.65, y+0.95, MUTED, arrow=True)
        y += 0.95

# ══════════════════════════════════════════════════════════════════════════════
#  INSPECTOR
# ══════════════════════════════════════════════════════════════════════════════

def inspect(path):
    """Heuristic check for overflowing text boxes and empty titles."""
    prs = Presentation(path)
    issues = []
    for si, slide in enumerate(prs.slides, 1):
        had_text = False
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            tf = shp.text_frame
            txt = "\n".join(p.text for p in tf.paragraphs)
            if txt.strip():
                had_text = True
            else:
                continue  # decorative / empty shape — not an overflow candidate
            # estimate overflow
            try:
                box_w_in = shp.width / EMU
                box_h_in = shp.height / EMU
            except Exception:
                continue
            if box_w_in <= 0 or box_h_in <= 0:
                continue
            if box_h_in < 0.2:
                continue  # thin accent bar / rule — ignore
            est_lines = 0
            max_pt = 0
            for p in tf.paragraphs:
                ptxt = p.text
                # font size for this paragraph
                sizes = [r.font.size.pt for r in p.runs if r.font.size is not None]
                fs = max(sizes) if sizes else 16
                max_pt = max(max_pt, fs)
                # chars per line at this font size in this box width
                char_w_in = fs * 0.0092            # empirical avg glyph advance
                cpl = max(1, int(box_w_in / char_w_in))
                est_lines += max(1, math.ceil(max(1, len(ptxt)) / cpl))
            needed_in = est_lines * (max_pt * 1.25) / 72.0 + 0.1
            if needed_in > box_h_in * 1.08:
                issues.append(f"  slide {si}: possible overflow — '{txt[:42].strip()}…' "
                              f"needs ~{needed_in:.2f}in, box {box_h_in:.2f}in")
        if not had_text:
            issues.append(f"  slide {si}: NO text found (blank?)")
    print(f"[inspect] {os.path.basename(path)} — {len(prs.slides._sldIdLst)} slides, "
          f"{len(issues)} issue(s)")
    for it in issues:
        print(it)
    return issues


if __name__ == "__main__":
    here = os.path.dirname(os.path.abspath(__file__))
    main_path = os.path.join(here, "OptiCart_Main.pptx")
    backup_path = os.path.join(here, "OptiCart_Backup.pptx")

    build_main(main_path)
    build_backup(backup_path)
    print("built:", main_path)
    print("built:", backup_path)

    i1 = inspect(main_path)
    i2 = inspect(backup_path)
    total = len(i1) + len(i2)
    print(f"\n[done] total heuristic issues: {total}")
